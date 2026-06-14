"""
PowerPoint (PPTX) generator for attribution reports.

Takes the same `summary` dict that drives the HTML and PDF, and produces a
native-editable 16:9 PowerPoint deck. Each section becomes one slide.

Design translation:
  HTML two-tone card        → stacked rectangles (tinted top + white bottom)
  Italic *word* accent      → italic text run in brand primary color
  Chart.js canvas           → PNG rendered by matplotlib, embedded as picture
  Instrument Serif titles   → embedded font if available, Georgia italic fallback

Slide size: 13.33 in × 7.5 in (16:9 widescreen standard).

Usage:
    from build_pptx import build_pptx_from_summary
    build_pptx_from_summary(summary, '/mnt/user-data/outputs/report.pptx')
"""
import io
import os
import sys
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from chart_renderer_png import render_chart_png


# =============================================================
# Slide geometry (16:9 widescreen)
# =============================================================
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Margins & grid
M_LEFT = Inches(0.55)
M_RIGHT = Inches(0.55)
M_TOP = Inches(0.55)
M_BOTTOM = Inches(0.55)
CONTENT_WIDTH = SLIDE_WIDTH - M_LEFT - M_RIGHT

# Fonts
FONT_DISPLAY = 'Instrument Serif'   # falls back to Georgia
FONT_DISPLAY_FALLBACK = 'Georgia'
FONT_BODY = 'Geist'
FONT_BODY_FALLBACK = 'Calibri'
FONT_MONO = 'Geist Mono'
FONT_MONO_FALLBACK = 'Consolas'



# ============================================================
# THEME SYSTEM
# ============================================================
# `theme` field on brand_profile: "light" (default) or "gradient" (dark bg).
# When gradient, we render a dark diagonal gradient background on every slide,
# flip text colors to white/light gray, and render cards as dark brand-tinted
# panels. Brand primary + accent stay bright so they pop on dark.

def _is_gradient_theme(brand):
    return (brand.get('theme') or '').lower() == 'gradient'


def _theme_colors(brand):
    """Returns theme-adapted color dict. Brand primary/accent stay the same."""
    if _is_gradient_theme(brand):
        return {
            'text_headline': '#FFFFFF',
            'text_body': '#F0F2F5',
            'text_muted': '#C8CDD7',
            'card_bg_top': '#1A2340',      # Dark mid-tone
            'card_bg_bot': '#0F1830',      # Darker bottom
            'card_divider': '#2F3A54',
            'card_border': '#2A3550',
            'chrome_rule': '#2F3A54',
            'chart_bg': '#0F1830',
            'chart_text': '#E6E8EE',
        }
    return {
        'text_headline': brand.get('colors', {}).get('ink', '#0B1B34'),
        'text_body': brand.get('colors', {}).get('ink', '#0B1B34'),
        'text_muted': brand.get('colors', {}).get('muted', '#64748B'),
        'card_bg_top': brand.get('colors', {}).get('surface', '#F5F7FB'),
        'card_bg_bot': '#FFFFFF',
        'card_divider': '#D5DBE3',
        'card_border': '#D5DBE3',
        'chrome_rule': '#E3E7ED',
        'chart_bg': '#FFFFFF',
        'chart_text': '#0B1B34',
    }


def _mix_hex(hex_a, hex_b, t):
    a = hex_a.lstrip("#"); b = hex_b.lstrip("#")
    r1, g1, b1 = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
    r2, g2, b2 = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
    r = int(r1 * (1 - t) + r2 * t)
    g = int(g1 * (1 - t) + g2 * t)
    bl = int(b1 * (1 - t) + b2 * t)
    return "#{:02X}{:02X}{:02X}".format(r, g, bl)


def _set_gradient_bg(slide, hex_start, hex_end, angle_deg=45):
    """Apply a linear gradient to the slide background via XML injection."""
    from lxml import etree
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    angle_pptx = int(angle_deg * 60000)
    start_hex = hex_start.lstrip("#").upper()
    end_hex = hex_end.lstrip("#").upper()
    bg_xml = (
        '<p:bg xmlns:p="' + p_ns + '" xmlns:a="' + a_ns + '">'
        '<p:bgPr>'
        '<a:gradFill rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="' + start_hex + '"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="' + end_hex + '"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="' + str(angle_pptx) + '" scaled="1"/>'
        '</a:gradFill>'
        '</p:bgPr>'
        '</p:bg>'
    )
    bg_elem = etree.fromstring(bg_xml)
    cSld = slide._element.find(qn("p:cSld"))
    for existing in cSld.findall(qn("p:bg")):
        cSld.remove(existing)
    cSld.insert(0, bg_elem)


def _apply_theme_background(slide, brand):
    """Apply gradient background if theme is gradient."""
    if not _is_gradient_theme(brand):
        return
    primary = brand.get("colors", {}).get("primary", "#0160BF")
    ink = brand.get("colors", {}).get("ink", "#0B1B34")
    # Start: deep ink, End: ink mixed with 35% primary (brand-tinted dark)
    end_hex = _mix_hex(ink, primary, 0.35)
    _set_gradient_bg(slide, ink, end_hex, angle_deg=45)


def _hex_to_rgb(hex_str, default='#0B1B34'):
    s = (hex_str or default).lstrip('#')
    if len(s) != 6:
        s = default.lstrip('#')
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _pale_rgb(hex_str, mix=0.92):
    """Return a pale (white-mixed) version of a color. mix=0.92 means 92% white."""
    s = (hex_str or '#0160BF').lstrip('#')
    if len(s) != 6:
        s = '0160BF'
    r = int(s[0:2], 16); g = int(s[2:4], 16); b = int(s[4:6], 16)
    pr = int(r * (1 - mix) + 255 * mix)
    pg = int(g * (1 - mix) + 255 * mix)
    pb = int(b * (1 - mix) + 255 * mix)
    return RGBColor(pr, pg, pb)



# =============================================================
# Text helpers
# =============================================================

def _add_text(text_frame, text, *, font=FONT_DISPLAY, size=18, color='#0B1B34',
              bold=False, italic=False, align=None, clear=True):
    """Add a paragraph with a single run. Returns the run for further styling."""
    if clear:
        text_frame.clear()
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_to_rgb(color)
    return run


def _add_title_with_italic_accent(text_frame, text, *, size=44, color='#0B1B34',
                                   accent_color='#0160BF', align=None, clear=True):
    """
    Parse a title like "Your *attribution* picture." — the *word* fragment becomes
    italic and brand-colored. Everything else is regular.
    """
    if clear:
        text_frame.clear()
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    if align:
        p.alignment = align

    segments = []
    in_italic = False
    buffer = []
    for ch in text:
        if ch == '*':
            if buffer:
                segments.append((''.join(buffer), in_italic))
                buffer = []
            in_italic = not in_italic
        else:
            buffer.append(ch)
    if buffer:
        segments.append((''.join(buffer), in_italic))

    for seg_text, seg_italic in segments:
        run = p.add_run()
        run.text = seg_text
        run.font.name = FONT_DISPLAY
        run.font.size = Pt(size)
        run.font.italic = seg_italic
        run.font.color.rgb = _hex_to_rgb(accent_color if seg_italic else color)


# =============================================================
# Shape primitives
# =============================================================

def _add_rect(slide, left, top, width, height, fill_hex, *, line=False):
    """Add a filled rectangle. Returns the shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    if line:
        shape.line.color.rgb = _hex_to_rgb('#0B1B34')
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_two_tone_card(slide, left, top, width, height, colors, *,
                        tint_fraction=0.45, brand=None):
    """
    Add a two-tone card: tinted top zone + bottom zone, hairline border.
    If `brand` has theme='gradient', uses dark-theme card colors instead.
    Returns (card_shape, top_height, bottom_height).
    """
    brand = brand or {}
    if _is_gradient_theme(brand):
        bg_top = '#1A2340'
        bg_bot = '#0F1830'
        border = '#2A3550'
        divider_color = '#2F3A54'
    else:
        bg_top = colors.get('surface', '#F5F7FB')
        bg_bot = '#FFFFFF'
        border = '#D5DBE3'
        divider_color = '#D5DBE3'

    top_h = Emu(int(height * tint_fraction))
    bot_h = height - top_h

    # Outer frame (bottom-zone color with hairline border)
    outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    outer.fill.solid()
    outer.fill.fore_color.rgb = _hex_to_rgb(bg_bot)
    outer.line.color.rgb = _hex_to_rgb(border)
    outer.line.width = Pt(0.6)
    outer.shadow.inherit = False

    # Top tinted zone
    top_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, top_h)
    top_zone.fill.solid()
    top_zone.fill.fore_color.rgb = _hex_to_rgb(bg_top)
    top_zone.line.fill.background()
    top_zone.shadow.inherit = False

    # Hairline divider between zones
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + top_h, width, Emu(6000))
    divider.fill.solid()
    divider.fill.fore_color.rgb = _hex_to_rgb(divider_color)
    divider.line.fill.background()
    divider.shadow.inherit = False

    return outer, top_h, bot_h


def _add_textbox(slide, left, top, width, height):
    """Add an empty textbox. Returns the shape."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return tb


# =============================================================
# Section renderers
# =============================================================

def _slide_header_and_footer(slide, brand, customer_domain, slide_num, total_slides,
                             skill_version, is_section_divider=False):
    """Add the report-wide header (logo + brand tag + brand-color accent rule)
    and footer (source + version + brand-color accent rule) on every slide.
    Section divider slides get different header treatment."""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    primary = colors.get('primary', '#0160BF')
    muted = colors.get('muted', '#8A93A0')
    ink = tc['text_headline']
    tc = _theme_colors(brand)

    # ---- Header band ----
    # Just the company name — no "× UTM Grabber" tag
    hdr = _add_textbox(slide, M_LEFT, Inches(0.22), Inches(6), Inches(0.35))
    p = hdr.text_frame.paragraphs[0]
    logo_run = p.add_run()
    logo_run.text = brand.get('company_name', 'UTM Grabber')
    logo_run.font.name = FONT_DISPLAY
    logo_run.font.size = Pt(17)
    logo_run.font.color.rgb = _hex_to_rgb(tc['text_headline'])

    # Small brand-color accent mark (a 20pt wide bar at full brand primary)
    # — visual weight in brand color, not a decorative stripe
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         M_LEFT, Inches(0.65),
                                         Inches(0.3), Emu(12000))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _hex_to_rgb(primary)
    accent_bar.line.fill.background()
    accent_bar.shadow.inherit = False

    # Faint horizontal rule under header, spanning most of slide
    hdr_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        M_LEFT + Inches(0.4), Inches(0.67),
                                        SLIDE_WIDTH - M_LEFT - M_RIGHT - Inches(0.4),
                                        Emu(6000))
    hdr_rule.fill.solid()
    hdr_rule.fill.fore_color.rgb = _hex_to_rgb(tc['chrome_rule'])
    hdr_rule.line.fill.background()
    hdr_rule.shadow.inherit = False

    # ---- Footer band ----
    # Faint horizontal rule above footer
    ftr_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        M_LEFT, SLIDE_HEIGHT - Inches(0.48),
                                        SLIDE_WIDTH - M_LEFT - M_RIGHT,
                                        Emu(6000))
    ftr_rule.fill.solid()
    ftr_rule.fill.fore_color.rgb = _hex_to_rgb(tc['chrome_rule'])
    ftr_rule.line.fill.background()
    ftr_rule.shadow.inherit = False

    # Small brand accent mark on footer too
    ftr_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         M_LEFT, SLIDE_HEIGHT - Inches(0.48),
                                         Inches(0.3), Emu(12000))
    ftr_accent.fill.solid()
    ftr_accent.fill.fore_color.rgb = _hex_to_rgb(primary)
    ftr_accent.line.fill.background()
    ftr_accent.shadow.inherit = False

    # Footer text (left)
    ftr_left = _add_textbox(slide, M_LEFT + Inches(0.4), SLIDE_HEIGHT - Inches(0.42),
                             Inches(7), Inches(0.25))
    _add_text(ftr_left.text_frame,
              f"Data source: UTM Grabber · {customer_domain}",
              font=FONT_MONO, size=10, color=tc['text_muted'])

    # Footer text (right)
    ftr_right = _add_textbox(slide, SLIDE_WIDTH - M_RIGHT - Inches(5),
                              SLIDE_HEIGHT - Inches(0.42),
                              Inches(5), Inches(0.25))
    today = datetime.now().strftime('%b %-d, %Y')
    ftr_text = f"v{skill_version} · {today}   ·   {slide_num} / {total_slides}"
    _add_text(ftr_right.text_frame, ftr_text,
              font=FONT_MONO, size=10, color=tc['text_muted'], align=PP_ALIGN.RIGHT)


def render_title_slide(slide, s, brand):
    """Section type: title-block"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    primary = colors.get('primary', '#0160BF')
    tc = _theme_colors(brand)

    # Kicker
    if s.get('kicker'):
        tb = _add_textbox(slide, M_LEFT, Inches(1.35), CONTENT_WIDTH, Inches(0.5))
        _add_text(tb.text_frame, s['kicker'].upper(),
                  font=FONT_MONO, size=17, color=primary)

    # Main title with italic accent
    tb = _add_textbox(slide, M_LEFT, Inches(2.0), CONTENT_WIDTH, Inches(3.2))
    _add_title_with_italic_accent(tb.text_frame, s.get('title', ''),
                                   size=76,
                                   color=tc['text_headline'],
                                   accent_color=primary)

    # Brand-primary decorative horizontal rule — short, right-aligned under title
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    M_LEFT, Inches(5.1),
                                    Inches(1.6), Emu(20000))
    deco.fill.solid()
    deco.fill.fore_color.rgb = _hex_to_rgb(primary)
    deco.line.fill.background()
    deco.shadow.inherit = False

    # Meta bits
    if s.get('meta_bits'):
        tb = _add_textbox(slide, M_LEFT, Inches(5.4), CONTENT_WIDTH, Inches(0.4))
        _add_text(tb.text_frame, '  ·  '.join(s['meta_bits']),
                  font=FONT_BODY, size=18, color=tc['text_muted'])


def render_stat_strip_slide(slide, s, brand):
    """Section type: stat-strip — grid of 2-4 two-tone cards, sized to fill the canvas"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    stats = s.get('stats', [])[:4]
    n = len(stats)
    if n == 0:
        return

    gap = Inches(0.22)
    card_w = Emu(int((CONTENT_WIDTH - gap * (n - 1)) / n))
    card_h = Inches(3.8)
    top = Inches(2.0)

    for i, stat in enumerate(stats):
        left = M_LEFT + (card_w + gap) * i
        _add_two_tone_card(slide, left, top, card_w, card_h, colors, tint_fraction=0.65, brand=brand)

        # Value (italic serif, big)
        value_tb = _add_textbox(slide, left + Inches(0.25), top + Inches(0.35),
                                card_w - Inches(0.5), Inches(2.3))
        _tc = _theme_colors(brand)
        _add_text(value_tb.text_frame, str(stat.get('value', '')),
                  font=FONT_DISPLAY, size=72,
                  color=_tc['text_headline'], italic=True)

        # Label + delta (bottom zone)
        label_tb = _add_textbox(slide, left + Inches(0.25), top + Inches(2.75),
                                card_w - Inches(0.5), Inches(0.9))
        label_tf = label_tb.text_frame
        label_tf.clear()
        p = label_tf.paragraphs[0]
        label_run = p.add_run()
        label_run.text = stat.get('label', '').upper()
        label_run.font.name = FONT_MONO
        label_run.font.size = Pt(13)
        label_run.font.color.rgb = _hex_to_rgb(colors.get('primary', '#0160BF'))

        if stat.get('delta_label'):
            direction = stat.get('delta_direction', 'flat')
            dcolor = {'up': '#1d7a4a', 'down': '#c4353a', 'flat': tc['text_muted']}.get(direction, tc['text_muted'])
            delta_p = label_tf.add_paragraph()
            delta_p.space_before = Pt(4)
            delta_run = delta_p.add_run()
            if direction == 'flat':
                delta_run.text = '—'
            elif direction == 'up':
                delta_run.text = f"▲ {stat['delta_label']}"
            else:
                delta_run.text = f"▼ {stat['delta_label']}"
            delta_run.font.name = FONT_MONO
            delta_run.font.size = Pt(13)
            delta_run.font.color.rgb = _hex_to_rgb(dcolor)


def render_hero_number_slide(slide, s, brand):
    """Section type: hero-number"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    card_w = Inches(10)
    card_h = Inches(5.2)
    left = (SLIDE_WIDTH - card_w) / 2
    top = Inches(1.2)

    _add_two_tone_card(slide, left, top, card_w, card_h, colors, tint_fraction=0.65, brand=brand)

    # Kicker
    if s.get('kicker'):
        kicker_tb = _add_textbox(slide, left + Inches(0.3), top + Inches(0.4),
                                  card_w - Inches(0.6), Inches(0.5))
        _add_text(kicker_tb.text_frame, s['kicker'].upper(),
                  font=FONT_MONO, size=13,
                  color=colors.get('primary', '#0160BF'),
                  align=PP_ALIGN.CENTER)

    # Value (enormous italic serif)
    value_tb = _add_textbox(slide, left + Inches(0.3), top + Inches(1.0),
                             card_w - Inches(0.6), Inches(2.5))
    _add_text(value_tb.text_frame, str(s.get('value', '')),
              font=FONT_DISPLAY, size=160,
              color=tc['text_headline'],
              italic=True, align=PP_ALIGN.CENTER)

    # Label (bottom zone)
    if s.get('label'):
        label_tb = _add_textbox(slide, left + Inches(0.6), top + Inches(3.7),
                                 card_w - Inches(1.2), Inches(1.3))
        _add_text(label_tb.text_frame, s['label'],
                  font=FONT_BODY, size=18,
                  color=tc['text_headline'],
                  align=PP_ALIGN.CENTER)


def render_chart_slide(slide, s, brand, chart_width_in=10.5, chart_height_in=5.0):
    """Section type: chart — full-width chart with caption"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    png_bytes = render_chart_png(s.get('chart', {}), colors,
                                   width_in=chart_width_in, height_in=chart_height_in)
    if not png_bytes:
        return

    # Chart container
    chart_w = Inches(chart_width_in)
    chart_h = Inches(chart_height_in)
    left = (SLIDE_WIDTH - chart_w) / 2
    top = Inches(1.1)
    _add_rect(slide, left, top, chart_w, chart_h, '#FFFFFF', line=True)
    slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=chart_w, height=chart_h)

    # Caption
    if s.get('caption'):
        cap_tb = _add_textbox(slide, M_LEFT, top + chart_h + Inches(0.15),
                              CONTENT_WIDTH, Inches(0.4))
        _add_text(cap_tb.text_frame, s['caption'],
                  font=FONT_MONO, size=12,
                  color=tc['text_muted'],
                  align=PP_ALIGN.CENTER)


def render_chart_insight_slide(slide, s, brand):
    """Section type: chart-insight — chart on left, two-tone insight card on right"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    # Chart area (left 60%)
    chart_w = Inches(7.5)
    chart_h = Inches(5.0)
    chart_left = M_LEFT
    chart_top = Inches(1.1)
    png_bytes = render_chart_png(s.get('chart', {}), colors,
                                   width_in=7.5, height_in=5.0)
    if png_bytes:
        _add_rect(slide, chart_left, chart_top, chart_w, chart_h, '#FFFFFF', line=True)
        slide.shapes.add_picture(io.BytesIO(png_bytes), chart_left, chart_top,
                                  width=chart_w, height=chart_h)

    if s.get('caption'):
        cap_tb = _add_textbox(slide, chart_left, chart_top + chart_h + Inches(0.12),
                              chart_w, Inches(0.4))
        _add_text(cap_tb.text_frame, s['caption'],
                  font=FONT_MONO, size=10,
                  color=tc['text_muted'],
                  align=PP_ALIGN.CENTER)

    # Insight card (right 40%)
    card_w = SLIDE_WIDTH - M_RIGHT - chart_left - chart_w - Inches(0.3)
    card_h = chart_h + Inches(0.5)
    card_left = chart_left + chart_w + Inches(0.3)
    card_top = chart_top
    _add_two_tone_card(slide, card_left, card_top, card_w, card_h, colors, tint_fraction=0.42, brand=brand)

    # Top: kicker + title
    if s.get('insight_kicker'):
        k_tb = _add_textbox(slide, card_left + Inches(0.25), card_top + Inches(0.25),
                             card_w - Inches(0.5), Inches(0.4))
        _add_text(k_tb.text_frame, s['insight_kicker'].upper(),
                  font=FONT_MONO, size=12,
                  color=colors.get('primary', '#0160BF'))

    if s.get('insight_title'):
        t_tb = _add_textbox(slide, card_left + Inches(0.25), card_top + Inches(0.65),
                             card_w - Inches(0.5), Inches(1.5))
        _add_title_with_italic_accent(t_tb.text_frame, s['insight_title'],
                                       size=26,
                                       color=tc['text_headline'],
                                       accent_color=colors.get('primary', '#0160BF'))

    # Bottom: body paragraphs
    if s.get('insight_body'):
        b_tb = _add_textbox(slide, card_left + Inches(0.25),
                             card_top + Inches(2.6),
                             card_w - Inches(0.5),
                             Inches(2.6))
        b_tf = b_tb.text_frame
        b_tf.clear()
        for i, para in enumerate(s['insight_body']):
            p = b_tf.paragraphs[0] if i == 0 else b_tf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = para
            run.font.name = FONT_BODY
            run.font.size = Pt(19)
            run.font.color.rgb = _hex_to_rgb(tc['text_body'])


def render_section_header_slide(slide, s, brand):
    """Section type: section-header — divider slide with kicker + big italic title + pale watermark numeral"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    primary = colors.get('primary', '#0160BF')
    ink = tc['text_headline']
    tc = _theme_colors(brand)

    # Number + kicker integrated on one line, brand primary
    # Like: "02 · SOURCES" where "02" is the section marker, "SOURCES" is the topic.
    k_tb = _add_textbox(slide, M_LEFT, Inches(3.0), Inches(11), Inches(0.5))
    k_tf = k_tb.text_frame
    k_tf.clear()
    p = k_tf.paragraphs[0]
    if s.get('number'):
        num_run = p.add_run()
        num_run.text = f"{s['number']}  "
        num_run.font.name = FONT_DISPLAY
        num_run.font.italic = True
        num_run.font.size = Pt(18)
        num_run.font.color.rgb = _hex_to_rgb(primary)

        sep_run = p.add_run()
        sep_run.text = '· '
        sep_run.font.name = FONT_MONO
        sep_run.font.size = Pt(13)
        sep_run.font.color.rgb = _hex_to_rgb(primary)

    if s.get('kicker'):
        kick_run = p.add_run()
        kick_run.text = s['kicker'].upper()
        kick_run.font.name = FONT_MONO
        kick_run.font.size = Pt(15)
        kick_run.font.color.rgb = _hex_to_rgb(primary)

    # Brand-color vertical accent line (short — just as a typographic ornament)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      M_LEFT, Inches(3.55),
                                      Inches(0.6), Emu(25000))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _hex_to_rgb(primary)
    accent.line.fill.background()
    accent.shadow.inherit = False

    # Big title with italic accent
    if s.get('title'):
        t_tb = _add_textbox(slide, M_LEFT, Inches(3.8), Inches(8), Inches(2.5))
        _add_title_with_italic_accent(t_tb.text_frame, s['title'],
                                       size=64, color=tc['text_headline'], accent_color=primary)


def render_recommendations_slide(slide, s, brand):
    """Section type: recommendations — 2x2 grid of two-tone cards"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    items = s.get('items', [])[:4]

    gap_x = Inches(0.25)
    gap_y = Inches(0.22)
    card_w = Emu(int((CONTENT_WIDTH - gap_x) / 2))
    card_h = Inches(3.0)
    start_top = Inches(0.92)

    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        left = M_LEFT + (card_w + gap_x) * col
        top = start_top + (card_h + gap_y) * row

        # Tint zone now covers 58% of card to fit 2-line titles cleanly
        _add_two_tone_card(slide, left, top, card_w, card_h, colors, tint_fraction=0.58, brand=brand)

        # Label (action number) — small, snug to top
        label = item.get('label', f'Action {str(i+1).zfill(2)}')
        lab_tb = _add_textbox(slide, left + Inches(0.28), top + Inches(0.22),
                               card_w - Inches(0.56), Inches(0.35))
        _add_text(lab_tb.text_frame, label.upper(),
                  font=FONT_MONO, size=11,
                  color=colors.get('primary', '#0160BF'))

        # Title with italic accent — room for 2 lines at 22pt
        if item.get('title'):
            t_tb = _add_textbox(slide, left + Inches(0.28), top + Inches(0.62),
                                 card_w - Inches(0.56), Inches(1.15))
            _add_title_with_italic_accent(t_tb.text_frame, item['title'],
                                           size=22,
                                           color=_theme_colors(brand)['text_headline'],
                                           accent_color=colors.get('primary', '#0160BF'))

        # Body — bumped to 17pt for readability
        if item.get('body'):
            b_tb = _add_textbox(slide, left + Inches(0.28), top + Inches(1.92),
                                 card_w - Inches(0.56), Inches(1.0))
            _add_text(b_tb.text_frame, item['body'],
                      font=FONT_BODY, size=17,
                      color=_theme_colors(brand)['text_body'])


def render_ranked_list_slide(slide, s, brand):
    """Section type: ranked-list — table with tinted header row"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    cols = s.get('columns', [])
    rows = s.get('rows', [])
    if not cols or not rows:
        return

    # Limit to top 8 rows to keep on one slide
    rows = rows[:8]

    table_left = M_LEFT
    table_top = Inches(1.1)
    table_width = CONTENT_WIDTH
    header_h = Inches(0.5)
    row_h = Inches(0.5)
    table_height = header_h + row_h * len(rows)

    # Outer border
    _add_rect(slide, table_left, table_top, table_width, table_height, '#FFFFFF', line=True)

    # Header row — tinted
    _add_rect(slide, table_left, table_top, table_width, header_h,
               colors.get('surface', '#F5F7FB'))

    # Column widths: distribute proportionally based on declared widths where possible
    n_cols = len(cols)
    col_width = Emu(int(table_width / n_cols))

    # Headers
    for i, c in enumerate(cols):
        cell_tb = _add_textbox(slide, table_left + col_width * i + Inches(0.12),
                                table_top + Inches(0.1),
                                col_width - Inches(0.24), header_h - Inches(0.2))
        align = {'right': PP_ALIGN.RIGHT, 'center': PP_ALIGN.CENTER}.get(
            c.get('align'), PP_ALIGN.LEFT)
        _add_text(cell_tb.text_frame, c.get('label', '').upper(),
                  font=FONT_MONO, size=11,
                  color=tc['text_muted'],
                  align=align)

    # Rows
    for ri, row in enumerate(rows):
        row_top = table_top + header_h + row_h * ri
        for ci, c in enumerate(cols):
            val = row.get(c.get('key'), '')
            cell_tb = _add_textbox(slide, table_left + col_width * ci + Inches(0.12),
                                    row_top + Inches(0.08),
                                    col_width - Inches(0.24), row_h - Inches(0.16))
            align = {'right': PP_ALIGN.RIGHT, 'center': PP_ALIGN.CENTER}.get(
                c.get('align'), PP_ALIGN.LEFT)
            ctype = c.get('type', 'name')

            if ctype == 'trend' and isinstance(val, dict):
                label = val.get('label', '')
                state = val.get('state', 'steady')
                color = {'steady': colors.get('primary', '#0160BF'),
                         'rising': '#1d7a4a', 'cooling': '#c97a1c',
                         'cold': '#c4353a'}.get(state, tc['text_muted'])
                _add_text(cell_tb.text_frame, label,
                          font=FONT_MONO, size=12, color=color, align=align)
            elif ctype in ('rank', 'number'):
                _add_text(cell_tb.text_frame, str(val),
                          font=FONT_DISPLAY, size=18,
                          color=tc['text_muted'] if ctype == 'rank' else tc['text_headline'],
                          align=align)
            elif ctype == 'source':
                _add_text(cell_tb.text_frame, str(val),
                          font=FONT_MONO, size=12,
                          color=tc['text_muted'], align=align)
            elif ctype == 'mono-number':
                _add_text(cell_tb.text_frame, str(val),
                          font=FONT_MONO, size=12,
                          color=tc['text_headline'], align=align)
            else:  # 'name' or default
                _add_text(cell_tb.text_frame, str(val),
                          font=FONT_BODY, size=13,
                          color=tc['text_headline'], align=align)

        # Row separator (except on last)
        if ri < len(rows) - 1:
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          table_left, row_top + row_h,
                                          table_width, Emu(3000))
            sep.fill.solid()
            sep.fill.fore_color.rgb = _hex_to_rgb('#EDEEF2')
            sep.line.fill.background()
            sep.shadow.inherit = False


def render_insight_card_slide(slide, s, brand):
    """Section type: insight-card — standalone two-tone callout"""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    card_w = Inches(11)
    card_h = Inches(4.5)
    left = (SLIDE_WIDTH - card_w) / 2
    top = Inches(1.5)

    _add_two_tone_card(slide, left, top, card_w, card_h, colors, tint_fraction=0.45, brand=brand)

    # Kicker
    if s.get('kicker'):
        k_tb = _add_textbox(slide, left + Inches(0.4), top + Inches(0.4),
                             card_w - Inches(0.8), Inches(0.5))
        _add_text(k_tb.text_frame, s['kicker'].upper(),
                  font=FONT_MONO, size=14,
                  color=colors.get('primary', '#0160BF'))

    # Title with italic accent
    if s.get('title'):
        t_tb = _add_textbox(slide, left + Inches(0.4), top + Inches(0.85),
                             card_w - Inches(0.8), Inches(1.2))
        _add_title_with_italic_accent(t_tb.text_frame, s['title'],
                                       size=32,
                                       color=tc['text_headline'],
                                       accent_color=colors.get('primary', '#0160BF'))

    # Body
    if s.get('body'):
        b_tb = _add_textbox(slide, left + Inches(0.4), top + Inches(2.3),
                             card_w - Inches(0.8), Inches(2.0))
        b_tf = b_tb.text_frame
        b_tf.clear()
        for i, para in enumerate(s['body']):
            p = b_tf.paragraphs[0] if i == 0 else b_tf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = para
            run.font.name = FONT_BODY
            run.font.size = Pt(17)
            run.font.color.rgb = _hex_to_rgb(tc['text_body'])


def render_closing_slide(slide, s, brand):
    """Section type: closing — compact summary card recapping the report's
    key numbers (not a CTA). Renders `summary_stats` (list of {value, label})
    and optional `bullets` (list of plain strings). Falls back to title+body
    for backward compat."""
    colors = brand.get('colors', {})
    tc = _theme_colors(brand)
    primary = colors.get('primary', '#0160BF')
    ink = tc['text_headline']
    muted = tc['text_muted']
    surface = colors.get('surface', '#F5F7FB')

    # Compact panel — sits high on the slide, leaves clearance for footer
    tc = _theme_colors(brand)
    panel_w = CONTENT_WIDTH
    panel_h = Inches(5.0)
    panel_left = M_LEFT
    panel_top = Inches(1.0)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     panel_left, panel_top, panel_w, panel_h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _hex_to_rgb(tc['card_bg_top'])
    panel.line.color.rgb = _hex_to_rgb(tc['card_border'])
    panel.line.width = Pt(0.6)
    panel.shadow.inherit = False

    # Kicker (brand-primary mono caps)
    kicker = (s.get('kicker') or 'At a glance').upper()
    k_tb = _add_textbox(slide, panel_left + Inches(0.7), panel_top + Inches(0.5),
                         panel_w - Inches(1.4), Inches(0.4))
    _add_text(k_tb.text_frame, kicker, font=FONT_MONO, size=13, color=primary)

    # Short brand-color rule under the kicker
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    panel_left + Inches(0.7),
                                    panel_top + Inches(0.9),
                                    Inches(0.8), Emu(22000))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _hex_to_rgb(primary)
    rule.line.fill.background()
    rule.shadow.inherit = False

    # Title (smaller — this isn't the hero anymore)
    if s.get('title'):
        t_tb = _add_textbox(slide, panel_left + Inches(0.7),
                             panel_top + Inches(1.15),
                             panel_w - Inches(1.4), Inches(1.0))
        _add_title_with_italic_accent(t_tb.text_frame, s['title'],
                                       size=40, color=tc['text_headline'], accent_color=primary)

    # Summary stats row
    stats = (s.get('summary_stats') or [])[:4]
    if stats:
        stats_top = panel_top + Inches(2.4)
        n = len(stats)
        gap = Inches(0.25)
        avail_w = panel_w - Inches(1.4)
        stat_w = Emu(int((avail_w - gap * (n - 1)) / n))
        for i, stat in enumerate(stats):
            left = panel_left + Inches(0.7) + (stat_w + gap) * i
            v_tb = _add_textbox(slide, left, stats_top, stat_w, Inches(0.95))
            _add_text(v_tb.text_frame, str(stat.get('value', '')),
                      font=FONT_DISPLAY, size=46, italic=True, color=tc['text_headline'])
            l_tb = _add_textbox(slide, left, stats_top + Inches(0.95),
                                 stat_w, Inches(0.35))
            _add_text(l_tb.text_frame, str(stat.get('label', '')).upper(),
                      font=FONT_MONO, size=10, color=tc['text_muted'])

    # Bullets (one-line each, bullet-dot in brand primary)
    bullets = (s.get('bullets') or [])[:3]
    if bullets:
        bullets_top = panel_top + Inches(3.9)
        for i, bullet in enumerate(bullets):
            b_tb = _add_textbox(slide, panel_left + Inches(0.7),
                                 bullets_top + Inches(0.32) * i,
                                 panel_w - Inches(1.4), Inches(0.32))
            b_tf = b_tb.text_frame
            b_tf.clear()
            p = b_tf.paragraphs[0]
            dot = p.add_run()
            dot.text = '·  '
            dot.font.name = FONT_MONO
            dot.font.size = Pt(14)
            dot.font.color.rgb = _hex_to_rgb(primary)
            r = p.add_run()
            r.text = str(bullet)
            r.font.name = FONT_BODY
            r.font.size = Pt(14)
            r.font.color.rgb = _hex_to_rgb(tc['text_body'])

    # Backward compat
    if not stats and not bullets and s.get('body'):
        b_tb = _add_textbox(slide, panel_left + Inches(0.7),
                             panel_top + Inches(2.5),
                             panel_w - Inches(1.4), Inches(1.0))
        _add_text(b_tb.text_frame, s['body'],
                  font=FONT_BODY, size=14, color=muted)


# =============================================================
# Main builder
# =============================================================

SECTION_RENDERERS = {
    'title-block': render_title_slide,
    'stat-strip': render_stat_strip_slide,
    'hero-number': render_hero_number_slide,
    'section-header': render_section_header_slide,
    'chart': render_chart_slide,
    'chart-insight': render_chart_insight_slide,
    'recommendations': render_recommendations_slide,
    'insight-card': render_insight_card_slide,
    'ranked-list': render_ranked_list_slide,
    'closing': render_closing_slide,
}


def build_pptx_from_summary(summary, output_path, validate=True):
    """
    Build a PowerPoint deck from a summary dict.
    One slide per section. Returns the output path.

    Set validate=False to skip schema validation (not recommended).
    """
    if validate:
        try:
            from validate_schema import validate_summary
            errors = validate_summary(summary)
            if errors:
                print(f"Warning: {len(errors)} schema issue(s) — proceeding anyway:")
                for e in errors[:5]:
                    print(f"  • {e}")
        except ImportError:
            pass
    meta = summary.get('meta', {})
    brand = meta.get('brand_profile', {}) or {}
    customer_domain = meta.get('customer_domain', '')
    skill_version = meta.get('skill_version', '0.9.0')

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # blank

    sections = summary.get('sections', [])
    total = len(sections)

    for idx, s in enumerate(sections):
        slide = prs.slides.add_slide(blank_layout)
        _apply_theme_background(slide, brand)
        stype = s.get('type')
        renderer = SECTION_RENDERERS.get(stype)
        if renderer:
            renderer(slide, s, brand)

        # Header + footer on every slide
        _slide_header_and_footer(slide, brand, customer_domain,
                                  idx + 1, total, skill_version)

    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == '__main__':
    import json
    if len(sys.argv) != 3:
        print("Usage: python build_pptx.py <summary.json> <output.pptx>")
        sys.exit(1)
    with open(sys.argv[1]) as f:
        summary = json.load(f)
    path = build_pptx_from_summary(summary, sys.argv[2])
    print(f"✓ Wrote PPTX: {path}")
